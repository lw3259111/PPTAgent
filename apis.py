import inspect
import os
import re
import traceback
from copy import deepcopy
from dataclasses import dataclass
from enum import Enum
from functools import partial

from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Pt

from llms import Role
from presentation import Picture, SlidePage


@dataclass
class HistoryMark:
    API_CALL_ERROR = "api_call_error"
    API_CALL_CORRECT = "api_call_correct"
    CODE_RUN_ERROR = "code_run_error"
    CODE_RUN_CORRECT = "code_run_correct"


# 让llm注意元素层级
# 想办法让模型不会同时编辑和删除
class CodeExecutor:

    def __init__(self, coder: Role, retry_times: int):
        self.api_history = []
        self.code_history = []
        self.coder = coder
        if coder is None:
            retry_times = 0
        self.retry_times = retry_times
        self.registered_functions = API_TYPES.all_funcs()
        self.function_regex = re.compile(r"^[a-z]+_[a-z]+\(.+\)$")

    def get_apis_docs(self, funcs: list[callable], show_example: bool = False):
        api_doc = []
        for func in funcs:
            sig = inspect.signature(func)
            params = []
            for name, param in sig.parameters.items():
                if name == "slide":
                    continue
                param_str = name
                if param.annotation != inspect.Parameter.empty:
                    param_str += f": {param.annotation.__name__}"
                if param.default != inspect.Parameter.empty:
                    param_str += f" = {repr(param.default)}"
                params.append(param_str)
            signature = f"def {func.__name__}({', '.join(params)})"
            if not show_example:
                api_doc.append(signature)
                continue
            doc = inspect.getdoc(func)
            api_doc.append(f"{signature}\n\t{doc}")
        return "\n".join(api_doc)

    def execute_actions(self, actions: str, edit_slide: SlidePage, error_time: int = 0):
        found_code = False
        api_calls = actions.strip().split("\n")
        backup_state = deepcopy(edit_slide)
        self.api_history.append(
            [HistoryMark.API_CALL_ERROR, edit_slide.slide_idx, actions]
        )
        for line_idx, line in enumerate(api_calls):
            try:
                if line_idx == len(api_calls) - 1 and not found_code:
                    raise ValueError("No code block found in the api call.")
                if line.startswith("def"):
                    raise ValueError("The function definition should not be output.")
                if not self.function_regex.match(line):
                    continue
                found_code = True
                func = line.split("(")[0]
                if func not in self.registered_functions:
                    raise ValueError(f"The function {func} is not defined.")
                self.code_history.append([HistoryMark.CODE_RUN_ERROR, line, None])
                partial_func = partial(self.registered_functions[func], edit_slide)
                eval(line, {}, {func: partial_func})
                self.code_history[-1][0] = HistoryMark.CODE_RUN_CORRECT
            except:
                error_time += 1
                trace_msg = traceback.format_exc()
                self.code_history[-1][-1] = trace_msg
                if error_time > self.retry_times:
                    return None
                api_lines = (
                    "\n".join(api_calls[: line_idx - 1])
                    + f"\n--> Error Line: {line}\n"
                    + "\n".join(api_calls[line_idx:])
                )
                actions = self.coder(
                    error_message=trace_msg,
                    faulty_api_sequence=api_lines,
                )
                return self.execute_actions(actions, backup_state, error_time)
        self.api_history[-1][0] = HistoryMark.API_CALL_CORRECT
        return edit_slide


def runs_merge(paragraph: _Paragraph):
    runs = paragraph.runs
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) == 1:
        return runs[0]
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def get_textframe(slide: SlidePage, textframe_id: str):
    if "_" not in textframe_id:
        raise ValueError("The textframe ID should contain a `_`, got: ", textframe_id)
    element_id, text_id = textframe_id.split("_")
    element_id, text_id = int(element_id), int(text_id)
    shape = slide.shapes[element_id]
    if not shape.text_frame.is_textframe:
        raise ValueError(f"The element {element_id} doesn't have a text frame.")
    for para in shape.text_frame.data:
        if para["idx"] == text_id:
            return shape, para
    raise ValueError(f"Incorrect textframe ID: {textframe_id}.")


def del_para(text: str, text_shape: BaseShape):
    for para in text_shape.text_frame.paragraphs:
        if para.text == text:
            para._element.getparent().remove(para._element)
            if len(text_shape.text_frame.paragraphs) == 0:
                text_shape.element.getparent().remove(text_shape.element)
            return
    raise ValueError(f"Incorrect shape: {text_shape}.")


# 这里有问题，因为可能遇到重复的文本，以后需要用id来区分(reverse一下)
def replace_para(orig_text: str, new_text: str, text_shape: BaseShape):
    for para in text_shape.text_frame.paragraphs:
        if para.text == orig_text:
            run = runs_merge(para)
            run.text = new_text
            return
    raise ValueError(f"Incorrect shape: {text_shape}.")


# 融合一下del textframe和del para
def del_textframe(slide: SlidePage, textframe_id: str):
    """Delete the textframe with the given id."""
    shape, para = get_textframe(slide, textframe_id)
    if textframe_id in shape.closures:
        raise ValueError(
            f"The textframe {textframe_id} has been edited, your should not delete it."
        )
    shape.closures[textframe_id] = partial(del_para, para["text"])


def replace_text(slide: SlidePage, textframe_id: str, text: str):
    """Replace the text of the textframe with the given id."""
    shape, para = get_textframe(slide, textframe_id)
    if textframe_id in shape.closures:
        raise ValueError(
            f"The textframe {textframe_id} has been edited, your should not edit it again."
        )
    shape.closures[textframe_id] = partial(replace_para, para["text"], text)


def set_font_style(
    slide: SlidePage,
    textframe_id: str,
    bold: bool = None,
    italic: bool = None,
    underline: bool = None,
    font_size: int = None,
    font_color: str = None,
):
    """
    Set the font style of a text frame, set the font color in Hexadecimal Color Notation.
    Example:
    >>> set_font_style("1_1", bold=True, font_size=24, font_color="FF0000")
    """
    shape, para = get_textframe(slide, textframe_id)
    paratext = para["text"]

    def set_font(text_shape: BaseShape):
        find = False
        if not text_shape.has_text_frame:
            raise ValueError(f"The element is not a text frame: {textframe_id}.")
        for para in shape.text_frame.paragraphs:
            if para.text == paratext:
                find = True
                break
        if not find:
            raise ValueError(f"Incorrect element id: {textframe_id}.")
        run = runs_merge(para)
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if underline is not None:
            run.font.underline = underline
        if font_size is not None:
            run.font.size = Pt(font_size)
        if font_color is not None:
            run.font.color.rgb = RGBColor.from_string(font_color)

    shape.closures[textframe_id] = set_font


def adjust_element_geometry(
    slide: SlidePage, element_id: str, left: int, top: int, width: int, height: int
):
    """
    Set the position and size of a element.

    Parameters:
    element_id (str, required): The ID of the element.
    left (int, required): The left position of the element.
    top (int, required): The top position of the element.
    width (int, required): The width of the element.
    height (int, required): The height of the element.

    Example:
    >>> set_shape_position("1", 100, 150, 200, 300)
    """
    shape = slide.shapes[int(element_id)]
    shape.left = Pt(left)
    shape.top = Pt(top)
    shape.width = Pt(width)
    shape.height = Pt(height)

    def set_geometry(shape: BaseShape):
        shape.left = left
        shape.top = top
        shape.width = width
        shape.height = height

    shape.closures[element_id] = set_geometry


def replace_image(slide: SlidePage, figure_id: str, image_path: str):
    """Replace the image of the element with the given id."""
    if not os.path.exists(image_path):
        raise ValueError(f"The image {image_path} does not exist.")
    shape = slide.shapes[int(figure_id)]
    if not isinstance(shape, Picture):
        raise ValueError("The element is not a Picture.")
    if figure_id in shape.closures:
        raise ValueError(
            f"The element {figure_id} has been edited, your should not edit it again."
        )
    shape.closures[figure_id] = lambda x: None
    shape.img_path = image_path


def del_element(slide: SlidePage, element_id: str):
    """Delete the element with the given id"""
    if "_" in element_id:
        raise ValueError(
            "Only the element_id of a textframe can contain a `_`, not an element."
        )
    shape = slide.shapes[int(element_id)]

    def del_shape(shape: BaseShape):
        shape.element.getparent().remove(shape.element)

    if shape.text_frame.is_textframe:
        for i in range(len(shape.text_frame.data)):
            if f"{element_id}_{i}" in shape.closures:
                raise ValueError(
                    f"The element {element_id} has been edited, your should not delete it."
                )
    if element_id in shape.closures:
        raise ValueError(
            f"The element {element_id} has been deleted, your should not delete it again."
        )
    shape.closures[element_id] = del_shape


class API_TYPES(Enum):
    Agent = [
        replace_text,
        del_textframe,
        replace_image,
        del_element,
    ]
    Coder = [
        replace_text,
        del_textframe,
        replace_image,
        del_element,
    ]
    Typographer = [
        set_font_style,
        adjust_element_geometry,
    ]

    # return all functions in the enum
    @classmethod
    def all_funcs(cls) -> dict[str, callable]:
        funcs = {}
        for attr in dir(cls):
            if attr.startswith("__"):
                continue
            funcs |= {func.__name__: func for func in getattr(cls, attr).value}
        return funcs
